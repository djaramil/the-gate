#!/usr/bin/env python3
"""Generate a professional PowerPoint presentation from the markdown presentation file."""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Paths
PRESENTATION_MD = Path("docs/PRESENTATION.md")
OUTPUT_PPTX = Path("docs/HootCamp_AI_Grading_Presentation.pptx")

# Professional color scheme
COLORS = {
    "primary": RGBColor(0x1A, 0x56, 0x7E),      # Deep blue
    "secondary": RGBColor(0x2D, 0x9C, 0xDB),    # Lighter blue
    "accent": RGBColor(0xF3, 0x9C, 0x12),       # Orange accent
    "dark": RGBColor(0x1E, 0x2A, 0x3A),         # Dark navy
    "text": RGBColor(0x2C, 0x3E, 0x50),         # Dark text
    "light_text": RGBColor(0xFF, 0xFF, 0xFF),   # White
    "code_bg": RGBColor(0xF5, 0xF5, 0xF5),      # Light gray
    "code_text": RGBColor(0x2C, 0x3E, 0x50),    # Code text
    "table_header": RGBColor(0x1A, 0x56, 0x7E), # Table header
    "table_alt": RGBColor(0xE8, 0xF4, 0xF8),    # Table alternate row
}


def parse_slides(md_path):
    """Parse markdown into slides."""
    content = md_path.read_text(encoding="utf-8")
    # Split by horizontal rules (---)
    slides_raw = re.split(r"\n---\s*\n", content)
    slides = []
    for raw in slides_raw:
        raw = raw.strip()
        if raw:
            slides.append(raw)
    return slides


def add_title_slide(prs, lines):
    """Add a title slide with title and subtitle."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add a professional background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["primary"]
    bg.line.fill.background()
    
    # Add accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.3), prs.slide_width, Inches(0.3))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["accent"]
    accent.line.fill.background()
    
    title = ""
    subtitle = ""
    for line in lines:
        line = line.strip()
        if line.startswith("# "):
            title = line[2:].strip()
        elif line.startswith("## "):
            subtitle = line[3:].strip()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.4), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS["light_text"]
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.4), Inches(1.0))
        tf = sub_box.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS["light_text"]
        p.alignment = PP_ALIGN.CENTER


def add_section_header_slide(prs, title):
    """Add a section header slide with a large title."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["dark"]
    bg.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.4), prs.slide_width, Inches(0.4))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["accent"]
    accent.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.4), Inches(2.0))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS["light_text"]
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, content_lines):
    """Add a standard content slide with title and bulleted content."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["primary"]
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS["light_text"]
    
    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(11.7), Inches(6.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for line in content_lines:
        line = line.strip()
        if not line:
            continue
        
        # Skip headings (already used as title)
        if line.startswith("#"):
            continue
        
        # Handle bullet points
        if line.startswith("- ") or line.startswith("* "):
            text = line[2:].strip()
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS["text"]
            p.level = 0
            p.space_after = Pt(10)
        elif line.startswith("  - ") or line.startswith("    - "):
            # Sub-bullets
            indent = len(line) - len(line.lstrip())
            text = line.strip()[2:].strip() if line.strip().startswith("- ") else line.strip()
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS["text"]
            p.level = 1 if indent <= 2 else 2
            p.space_after = Pt(6)
        elif line.startswith("1. ") or line.startswith("2. ") or line.startswith("3. ") or line.startswith("4. ") or line.startswith("5. ") or line.startswith("6. ") or line.startswith("7. ") or line.startswith("8. ") or line.startswith("9. "):
            # Numbered list
            text = line[3:].strip()
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS["text"]
            p.level = 0
            p.space_after = Pt(10)
        elif line.startswith("**") and line.endswith("**"):
            # Bold standalone
            p = tf.add_paragraph()
            p.text = line.strip("*")
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = COLORS["primary"]
            p.space_after = Pt(12)
        else:
            p = tf.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS["text"]
            p.space_after = Pt(8)


def add_table_slide(prs, title, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["primary"]
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS["light_text"]
    
    if not rows:
        return
    
    # Determine columns from first row
    cols = len(rows[0])
    if cols == 0:
        return
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.0)
    height = Inches(5.5)
    
    table = slide.shapes.add_table(len(rows), cols, left, top, width, height).table
    
    # Set column widths
    col_width = width / cols
    for col in table.columns:
        col.width = int(col_width)
    
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = cell_text.strip()
            
            # Format header row
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["table_header"]
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.color.rgb = COLORS["light_text"]
                    p.alignment = PP_ALIGN.CENTER
            else:
                # Alternate row colors
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS["table_alt"]
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(14)
                    p.font.color.rgb = COLORS["text"]
                    p.alignment = PP_ALIGN.CENTER


def add_code_slide(prs, title, code_lines):
    """Add a slide with code block."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["primary"]
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS["light_text"]
    
    # Code background
    code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.0), Inches(6.0))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = COLORS["code_bg"]
    code_bg.line.color.rgb = COLORS["secondary"]
    code_bg.line.width = Pt(2)
    
    # Code text
    code_text = "\n".join(code_lines)
    code_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(11.6), Inches(5.6))
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.name = "Courier New"
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS["code_text"]


def add_diagram_slide(prs, title, diagram_lines):
    """Add a slide with ASCII diagram."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["primary"]
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS["light_text"]
    
    # Diagram background
    diag_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.4), Inches(12.4), Inches(6.2))
    diag_bg.fill.solid()
    diag_bg.fill.fore_color.rgb = COLORS["code_bg"]
    diag_bg.line.color.rgb = COLORS["secondary"]
    diag_bg.line.width = Pt(2)
    
    # Diagram text
    diagram_text = "\n".join(diagram_lines)
    # Use smaller font for diagrams
    font_size = Pt(9) if len(diagram_lines) > 25 else Pt(10)
    
    diag_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(12.0), Inches(5.9))
    tf = diag_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = diagram_text
    p.font.name = "Courier New"
    p.font.size = font_size
    p.font.color.rgb = COLORS["code_text"]


def process_table(content_lines):
    """Extract table rows from markdown."""
    rows = []
    for line in content_lines:
        if line.strip().startswith("|") and line.strip().endswith("|"):
            # Skip separator lines
            if "---" in line:
                continue
            cells = [c.strip() for c in line.strip("|").split("|")]
            if cells and any(cells):
                rows.append(cells)
    return rows


def process_slide(prs, slide_text):
    """Process a single markdown slide and create the appropriate PPTX slide."""
    lines = slide_text.split("\n")
    
    # Clean empty lines from beginning and end
    while lines and not lines[0].strip():
        lines.pop(0)
    while lines and not lines[-1].strip():
        lines.pop()
    
    if not lines:
        return
    
    # Determine title
    title = ""
    content_start = 0
    for i, line in enumerate(lines):
        stripped = line.strip()
        if stripped.startswith("# "):
            title = stripped[2:]
            content_start = i + 1
            break
        elif stripped.startswith("## "):
            title = stripped[3:]
            content_start = i + 1
            break
    
    # If first line is just title with no other content, treat as section header
    content_lines = lines[content_start:]
    
    # Check if this is a title slide (has # title and ## subtitle)
    if title and content_start > 0 and not any(l.startswith("#") for l in content_lines):
        # Could be title slide or regular content
        has_subtitle = False
        for line in lines:
            if line.strip().startswith("## "):
                has_subtitle = True
                break
        if has_subtitle and title and not any(l.startswith("###") or l.startswith("- ") for l in content_lines):
            add_title_slide(prs, lines)
            return
    
    # If title is empty, use first line as title
    if not title:
        title = lines[0].strip().strip("#").strip()
        content_lines = lines[1:]
    
    # Section header slide if no content
    if not content_lines:
        add_section_header_slide(prs, title)
        return
    
    # Check for table
    table_lines = [l for l in content_lines if l.strip().startswith("|")]
    if table_lines and len(table_lines) >= 2:
        # Check if table is a proper markdown table
        if all("|" in l for l in table_lines):
            rows = process_table(table_lines)
            if rows:
                add_table_slide(prs, title, rows)
                return
    
    # Check for code block
    if "```" in slide_text:
        # Find code block
        code_match = re.search(r"```(?:\w*\n)?(.*?)```", slide_text, re.DOTALL)
        if code_match:
            code = code_match.group(1).strip().split("\n")
            # If the slide is mostly code, use code slide
            non_code = [l for l in content_lines if l.strip() and not l.strip().startswith("```")]
            if len(code) > 3 and len(non_code) <= 5:
                add_code_slide(prs, title, code)
                return
    
    # Check for diagram (contains lots of box drawing chars or looks like ASCII art)
    diagram_chars = sum(1 for line in content_lines if any(c in line for c in ["┌", "┐", "└", "┘", "│", "─", "├", "┤", "┬", "┴", "┼"]))
    if diagram_chars >= 5:
        diagram_lines = [l for l in content_lines if l.strip()]
        if diagram_lines:
            add_diagram_slide(prs, title, diagram_lines)
            return
    
    # Standard content slide
    add_content_slide(prs, title, content_lines)


def main():
    """Generate the professional PPTX presentation."""
    print(f"Reading markdown from {PRESENTATION_MD}")
    slides = parse_slides(PRESENTATION_MD)
    print(f"Found {len(slides)} slides")
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for i, slide_text in enumerate(slides):
        try:
            process_slide(prs, slide_text)
            print(f"  Slide {i+1} created")
        except Exception as e:
            print(f"  Error on slide {i+1}: {e}")
            # Add a blank slide with error text
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
            box.text_frame.text = f"Error processing slide {i+1}"
    
    prs.save(OUTPUT_PPTX)
    print(f"Saved presentation to {OUTPUT_PPTX}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
